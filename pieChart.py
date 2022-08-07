#create PIE CHART

import openpyxl
from pptx import Presentation
from openpyxl.chart import BarChart3D,AreaChart, Reference
import pandas as pd
import matplotlib.pyplot as plt
from pptx.util import Inches, Cm, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
import matplotlib.ticker as mtick
from PIL import Image
import numpy
from commonFunctions import getTitle, getSlideNo, getMonthlyData, px_to_inches, findDimensions, fixData

#for making line charts
def makePieChart(dataframe, prs, slideName, thisTitle, slideNo):

    #extract the subset of Data with values of Months for each segment and transponse for correct
    #x and y axes on area chart
    df2 = getMonthlyData(dataframe, 6, 19)
    df2[7:] = df2[7:].astype(int)
    nan_value = float("NaN")
    df2.replace("", nan_value, inplace=True)
    df2.dropna(how='all', axis=1, inplace=True)
    df2 = df2.T

    #choose releavnt data values
    df2.columns = df2.iloc[0]
    length = len(df2.columns)
    df2 = df2[1:].T
    df2 = df2.iloc[:length-2, :]

    #take percentage to scale y axis to be in range of 0 to 100
    #df2 = df2.apply(lambda x: x * 100 / sum(x), axis=1)

    #function to turn percentage in pie chart to absolute values
    def absolute_value(val):
        a = numpy.round(val / 100. * df2.iloc[:, 0].sum(), 0)
        return int(a)

    # create a pie chart
    ax = df2.plot.pie(subplots=True, autopct=absolute_value)

    # set the title and fix legend position to be at bottom center
    plt.title(thisTitle)
    plt.legend(loc='lower center', ncol = length, bbox_to_anchor=(0.6, -0.1), fontsize= 14)
    plt.legend.include_in_layout = False
    # show the plot and save it as a figure to put into ppt
    fig = plt.savefig('pieChart.jpg', dpi=1200, transparent=True, bbox_inches = 'tight')
    img = 'pieChart.jpg'

    #create new slide to add to presentation and set it to widescreen
    slide = prs.slides[slideNo - 1]
    slide.name = slideName
    # adjust title shape and position
    titleCurr = slide.shapes.title
    titleCurr.text = slideName
    titleCurr.text_frame.paragraphs[0].font.size = Pt(26)
    left = top = Inches(0.1)
    height = Inches(0.7)
    width = Inches(3)
    titleCurr.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    #center image with correct dimensions
    # imgD = px_to_inches(img)
    # dims = findDimensions(slide_size, imgD)
    txBox = slide.shapes.add_textbox(Inches(12.5), Inches(6.5), 7, 7)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = str(slideNo)
    p.font.size = Pt(20)
    #add picture and save to ppt
    # pic = slide.shapes.add_picture(img, Inches(0.5), Inches(1.6), Inches(12.15), Inches(5.33))
    pic = slide.shapes.add_picture(img, Cm(8.0), Cm(4.0), Cm(15), Cm(13))





