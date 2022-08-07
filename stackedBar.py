# to make STACKED BAR CHART

import openpyxl
from pptx import Presentation
from openpyxl.chart import BarChart3D, AreaChart, Reference
import pandas as pd
import matplotlib.pyplot as plt
from pptx.util import Inches, Cm, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
import matplotlib.ticker as mtick
from PIL import Image
from commonFunctions import getTitle, getSlideNo, getMonthlyData, px_to_inches, findDimensions, fixData


def makeStackedBarChart(dataframe, prs, slideName, thisTitle, slideNo, yaxis=''):
    # extract the subset of Data with values of Months for each segment and transponse for correct
    # x and y axes on area chart
    df2 = getMonthlyData(dataframe, 6, 19)
    nan_value = float("NaN")
    df2.replace("", nan_value, inplace=True)
    df2.dropna(how='all', axis=1, inplace=True)
    df2 = df2.T

    # choose releavnt data values
    df2.columns = df2.iloc[0]
    df2 = df2[1:]

    # colors = {'>30 days': 'blue', '21-30': 'red', 'same day': 'green', '11-20 days': 'red', '1-10 days': 'yellow'}

    # create a bar chart (stacked) with labels
    ax = df2.plot.barh(stacked=True)
    for bars in ax.containers:
        ax.bar_label(bars, label_type='center', fontsize=10)


    # set the title and fix legend position to be at bottom center
    plt.title(thisTitle)
    plt.legend(loc='lower center', ncol=len(df2.columns), bbox_to_anchor=(0.5, -0.3), fontsize='xx-small')
    plt.legend.include_in_layout = False
    # plt.ylabel(yaxis)
    # show the plot and save it as a figure to put into ppt
    fig = plt.savefig('stackedBarChart.jpg', dpi=1200, transparent=True, bbox_inches = 'tight')
    img = 'stackedBarChart.jpg'

    # create new slide to add to presentation and set it to widescreen
    # slide_size = (16, 9)

    # add slide of layout 'blank' ACCORDING TO company template and add it to presentation
    # slide_layout = prs.slide_layouts[0]
    # slide = prs.slides.add_slide(slide_layout)
    slide = prs.slides[slideNo - 1]

    # name the slide and add a title for the slide
    slide.name = slideName
    #title_placeholder = slide.shapes.title
    #title_placeholder.text = thisTitle
    # adjust title shape and position
    titleCurr = slide.shapes.title
    titleCurr.text = slideName
    titleCurr.text_frame.paragraphs[0].font.size = Pt(26)
    left = top = Inches(0.1)
    height = Inches(0.7)
    width = Inches(3)
    titleCurr.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    # shapes = slide.shapes
    # left = top = Inches(0.1)
    # height = Inches(0.7)
    # width = Inches(3)
    # shape = shapes.add_shape(
    #     MSO_SHAPE.RECTANGLE, left, top, width, height
    # )
    # fill = shape.fill
    # fill.solid()
    # fill.fore_color.rgb = RGBColor(255, 255, 255)
    # line = shape.line
    # line.color.rgb = RGBColor(255, 255, 255)
    # shape.text = slideName
    # text_frame = shape.text_frame
    # text_frame.text = slideName
    # p = text_frame.paragraphs[0]
    # font = p.font
    # font.color.rgb = RGBColor(0x00, 0x00, 0x00)
    # font.size = Pt(26)
    # text_frame.margin_bottom = Inches(0.08)
    # text_frame.margin_left = 0
    # text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    # text_frame.word_wrap = False
    # text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

    # center image with correct dimensions
    # imgD = px_to_inches(img)
    # dims = findDimensions(slide_size, imgD)
    txBox = slide.shapes.add_textbox(Inches(12.5), Inches(6.5), 7, 7)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = str(slideNo)
    p.font.size = Pt(20)
    # add picture and save to ppt
    # pic = slide.shapes.add_picture(img, Inches(0.5), Inches(1.6), Inches(12.15), Inches(5.33))
    pic = slide.shapes.add_picture(img, Cm(2.75), Cm(5.0), Cm(27.54), Cm(12.08))