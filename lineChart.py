# to make LINE CHART

import openpyxl
from pptx import Presentation
from openpyxl.chart import BarChart3D, AreaChart, Reference
import pandas as pd
import matplotlib.pyplot as plt
from pptx.util import Inches
from pptx.util import Inches, Cm, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
import matplotlib.ticker as mtick
from PIL import Image
from commonFunctions import add_value_labels, getTitle, getSlideNo, getMonthlyData, px_to_inches, findDimensions, fixData


def makeLineChart(dataframe, prs, slideName, thisTitle, slideNo, yaxis=''):

    # extract the subset of Data with values of Months for each segment and transponse for correct
    # x and y axes on area chart
    df2 = getMonthlyData(dataframe, 6, 20)
    nan_value = float("NaN")
    df2.replace("", nan_value, inplace=True)
    df2.dropna(how='all', axis=1, inplace=True)
    df2 = df2.T

    # choose relevant data values
    df2.columns = df2.iloc[0]
    df2 = df2[1:]

    if(thisTitle=="Target adherence"):
        print(df2)
    # create an line plot
    ax = df2.plot.line()

    # set the title and fix legend position to be at bottom center
    plt.title(thisTitle)
    plt.legend(loc='lower center', ncol=len(df2.columns), bbox_to_anchor=(0.5, -0.15), fontsize='xx-small')
    # show the plot and save it as a figure to put into ppt
    plt.legend.include_in_layout = False
    # plt.ylabel(yaxis)
    add_value_labels(ax, typ='line')
    fig = plt.savefig('lineChart.jpg', transparent=True, dpi=1200, bbox_inches = 'tight')
    img = 'lineChart.jpg'

    # create new slide to add to presentation and set it to widescreen
    slide = prs.slides[slideNo - 1]

    # name the slide and add a title for the slide
    slide.name = slideName

    # adjust title shape and position
    titleCurr = slide.shapes.title
    titleCurr.text = slideName
    titleCurr.text_frame.paragraphs[0].font.size = Pt(26)
    left = top = Inches(0.1)
    height = Inches(0.7)
    width = Inches(3)
    titleCurr.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    txBox = slide.shapes.add_textbox(Inches(12.5), Inches(6.5), 7, 7)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = str(slideNo)
    p.font.size = Pt(20)
    # add picture and save to ppt
    pic = slide.shapes.add_picture(img, Cm(2.75), Cm(5.0), Cm(27.54), Cm(12.08))