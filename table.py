import numpy as np
import openpyxl
from pptx import Presentation
from openpyxl.chart import BarChart3D,AreaChart, Reference
import pandas as pd
import matplotlib.pyplot as plt
from pptx.util import Inches, Cm, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
import matplotlib.ticker as mtick
from PIL import Image
from math import trunc
from PandasToPowerpoint.pd2ppt.pd2ppt import df_to_table
from commonFunctions import getTitle, isfloat, getSlideNo, getMonthlyData, px_to_inches, findDimensions, fixData

def makeTablePlot(df, prs, slideName, thisTitle, slideNo):


    # create new slide to add to presentation and set it to widescreen
    slide_size = (16, 9)
    slide = prs.slides[slideNo - 1]
    slide.name = slideName
    # adjust title shape and position
    titleCurr = slide.shapes.title
    titleCurr.text = slideName
    titleCurr.text_frame.paragraphs[0].font.size = Pt(26)
    left = top = Inches(0.1)
    height = Inches(0.7)
    width = Inches(3)
    titleCurr.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    df2 = getMonthlyData(df, 6, 20)
    nan_value = float("NaN")
    df2.replace("", nan_value, inplace=True)
    df2.dropna(how='all', axis=1, inplace=True)

    df_to_table(slide, df2, left=Inches(0.5), top=Inches(1.6), width=Inches(12.15), height=Inches(5.33))
    txBox = slide.shapes.add_textbox(Inches(12.5), Inches(6.5), 7, 7)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = str(slideNo)
    p.font.size = Pt(20)
    # print(thisTitle)
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            for cell in table.iter_cells():
                if isfloat(cell.text):
                    cell.text = str(trunc(float(cell.text)))
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER
                    for run in paragraph.runs:
                        run.font.size = Pt(17)

            for col in table.columns:
                lenCol = len(table.columns)
                wid = 7.43
                wid2 = wid / lenCol
                col.width = Inches(wid2)
            table.columns[0].width = Cm(15)
