# to make STACKED AREA CHARTS

import openpyxl
from pptx import Presentation
from openpyxl.chart import BarChart3D, AreaChart, Reference
import pandas as pd
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
import matplotlib.pyplot as plt
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Cm, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_COLOR_TYPE
import matplotlib.ticker as mtick
from PIL import Image
from commonFunctions import getTitle, getSlideNo, getMonthlyData, px_to_inches, findDimensions, fixData


def makeAreaChart(dataframe, prs, slideName, thisTitle, slideNo, yaxis=''):

    # extract the subset of Data with values of Months for each segment and transponse for correct
    # x and y axes on area chart
    df2 = getMonthlyData(dataframe, 6, 20)
    nan_value = float("NaN")
    df2.replace("", nan_value, inplace=True)
    df2.dropna(how='all', axis=1, inplace=True)
    df2 = df2.T
    # print(df2)

    # choose releavnt data values
    df2.columns = df2.iloc[0]
    df2 = df2[1:]

    # create an area plot (stacked)
    ax = df2.plot.area(stacked=True)

    # take percentage to scale y axis to be in range of 0 to 100
    # df2 = df2.apply(lambda x: x * 100 / sum(x), axis=1)

    # set the title and fix legend position to be at bottom center
    plt.title(thisTitle)
    plt.legend(loc='lower center', ncol=len(df2.columns), bbox_to_anchor=(0.5, -0.15), fontsize='xx-small')
    plt.legend.include_in_layout = False
    # plt.ylabel(yaxis)

    # show the plot and save it as a figure to put into ppt
    fig = plt.savefig('areaChart.jpg', transparent=True, dpi=1200, bbox_inches='tight')
    img = 'areaChart.jpg'
    slide = prs.slides[slideNo - 1]

    # name the slide and add a title for the slide
    slide.name = slideName
    titleCurr = slide.shapes.title
    titleCurr.text = slideName
    titleCurr.text_frame.paragraphs[0].font.size = Pt(26)
    left = top = Inches(0.1)
    height = Inches(0.7)
    width = Inches(3)
    titleCurr.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    txBox = slide.shapes.add_textbox(Inches(12.5), Inches(6.5), 7, 7)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = str(slideNo)
    p.font.size = Pt(20)

    pic = slide.shapes.add_picture(img, Cm(2.75), Cm(5.0), Cm(27.54), Cm(12.08))